import os
os.environ["STREAMLIT_SERVER_FILE_WATCHER_TYPE"] = "none"

import streamlit as st
import cv2
import numpy as np
from pptx import Presentation
import io
import gc

def fix_drawn_box_in_image(image_bytes):
    try:
        file_bytes = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(file_bytes, cv2.IMREAD_COLOR)
        
        if img is None:
            return image_bytes, False

        hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # 1. Color Masking (Green Marker & Dark Black Hand-Strokes)
        # Green Marker Range
        lower_green = np.array([30, 40, 40])
        upper_green = np.array([85, 255, 255])
        green_mask = cv2.inRange(hsv, lower_green, upper_green)

        # Dark Black Stroke Range
        lower_black = np.array([0, 0, 0])
        upper_black = np.array([180, 255, 40])
        black_mask = cv2.inRange(hsv, lower_black, upper_black)

        # Combine Both Markers
        marker_mask = cv2.bitwise_or(green_mask, black_mask)

        # Noise Filter (PPT layout ke baki lines ko ignore karne ke liye)
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (5, 5))
        marker_mask = cv2.morphologyEx(marker_mask, cv2.MORPH_OPEN, kernel)
        marker_mask = cv2.dilate(marker_mask, kernel, iterations=2)

        # 2. Find Marker Contours
        contours, _ = cv2.findContours(marker_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

        best_rect = None
        max_area = 0

        for cnt in contours:
            area = cv2.contourArea(cnt)
            # Area limit to isolate hand-drawn shape inside slide
            if 400 < area < (img.shape[0] * img.shape[1] * 0.35):
                if area > max_area:
                    max_area = area
                    best_rect = cv2.boundingRect(cnt)

        if best_rect is not None:
            x, y, w, h = best_rect

            # 3. Create Specific Mask for Inpainting
            inpaint_mask = np.zeros_like(gray)
            cv2.rectangle(inpaint_mask, (x, y), (x + w, y + h), 255, -1)
            inpaint_mask = cv2.bitwise_and(inpaint_mask, marker_mask)
            inpaint_mask = cv2.dilate(inpaint_mask, kernel, iterations=3)

            # 4. Erase Tedhi Line (Inpaint Original Background)
            clean_img = cv2.inpaint(img, inpaint_mask, inpaintRadius=7, flags=cv2.INPAINT_TELEA)

            # 5. Draw Perfect Straight Green Rectangle Box
            cv2.rectangle(clean_img, (x, y), (x + w, y + h), (0, 255, 0), 4)

            _, encoded_img = cv2.imencode('.png', clean_img)
            
            del img, hsv, gray, green_mask, black_mask, marker_mask, clean_img
            return encoded_img.tobytes(), True

    except Exception:
        pass

    return image_bytes, False

def process_ppt(ppt_bytes):
    prs = Presentation(io.BytesIO(ppt_bytes))
    modified = False

    total_slides = len(prs.slides)
    progress_bar = st.progress(0)

    for idx, slide in enumerate(prs.slides):
        shapes_to_replace = [s for s in slide.shapes if s.shape_type == 13]
        
        for shape in shapes_to_replace:
            try:
                image_bytes = shape.image.blob
                fixed_image_bytes, is_modified = fix_drawn_box_in_image(image_bytes)
                
                if is_modified:
                    modified = True
                    shape.image._blob = fixed_image_bytes
                
                del image_bytes
            except Exception:
                continue
        
        progress_bar.progress(int((idx + 1) / total_slides * 100))
        gc.collect()

    out_stream = io.BytesIO()
    prs.save(out_stream)
    out_stream.seek(0)
    
    progress_bar.empty()
    return out_stream.getvalue(), modified

# --- STREAMLIT UI ---
st.set_page_config(page_title="PPT Recce Mark Corrector", layout="centered")

st.title("PPT Recce Mark Corrector")
st.write("PPT upload karein, tool hand-drawn markers erase karke straight green box draw kar dega.")

uploaded_ppt = st.file_uploader("Apni PowerPoint (.pptx) file choose karein", type=["pptx"])

if uploaded_ppt is not None:
    if st.button("Process PPT", type="primary"):
        with st.spinner("Processing Slides... Please wait..."):
            processed_ppt_bytes, status = process_ppt(uploaded_ppt.read())
            
            if status:
                st.success("PPT Successfully Processed!")
                st.download_button(
                    label="Download Fixed PPT",
                    data=processed_ppt_bytes,
                    file_name=f"Fixed_{uploaded_ppt.name}",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.warning("PPT ke images me koi marker box detect nahi hua.")
